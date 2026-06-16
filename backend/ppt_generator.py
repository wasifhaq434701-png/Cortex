import io
import re
import json
import base64
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend for server-side rendering
import matplotlib.pyplot as plt
from matplotlib.figure import Figure

# ==========================================
# PPT Theme Definitions
# ==========================================
PPT_THEMES = {
    "dark": {
        "name": "Professional Dark",
        "bg": RGBColor(15, 23, 42),
        "accent": RGBColor(56, 189, 248),
        "text": RGBColor(226, 232, 240),
        "text_secondary": RGBColor(148, 163, 184),
        "surface": RGBColor(30, 41, 59),
        "chart_bg": "#0f172a",
        "chart_text": "#e2e8f0",
        "chart_colors": ['#38bdf8', '#a855f7', '#ec4899', '#10b981', '#f59e0b', '#ef4444'],
        "chart_grid": '#334155',
    },
    "corporate": {
        "name": "Corporate Blue",
        "bg": RGBColor(255, 255, 255),
        "accent": RGBColor(30, 64, 175),
        "text": RGBColor(30, 41, 59),
        "text_secondary": RGBColor(100, 116, 139),
        "surface": RGBColor(241, 245, 249),
        "chart_bg": "#ffffff",
        "chart_text": "#1e293b",
        "chart_colors": ['#1e40af', '#3b82f6', '#60a5fa', '#93c5fd', '#2563eb', '#1d4ed8'],
        "chart_grid": '#e2e8f0',
    },
    "minimal": {
        "name": "Minimal Light",
        "bg": RGBColor(250, 250, 250),
        "accent": RGBColor(15, 23, 42),
        "text": RGBColor(30, 41, 59),
        "text_secondary": RGBColor(100, 116, 139),
        "surface": RGBColor(255, 255, 255),
        "chart_bg": "#fafafa",
        "chart_text": "#1e293b",
        "chart_colors": ['#0f172a', '#475569', '#94a3b8', '#cbd5e1', '#334155', '#64748b'],
        "chart_grid": '#e2e8f0',
    },
    "neon": {
        "name": "Neon Cyberpunk",
        "bg": RGBColor(10, 0, 20),
        "accent": RGBColor(0, 255, 136),
        "text": RGBColor(224, 255, 240),
        "text_secondary": RGBColor(0, 200, 120),
        "surface": RGBColor(20, 10, 40),
        "chart_bg": "#0a0014",
        "chart_text": "#e0fff0",
        "chart_colors": ['#00ff88', '#ff00ff', '#00d4ff', '#ffff00', '#ff6600', '#ff0066'],
        "chart_grid": '#1a0a2e',
    },
    "sunset": {
        "name": "Warm Sunset",
        "bg": RGBColor(30, 15, 20),
        "accent": RGBColor(251, 146, 60),
        "text": RGBColor(255, 237, 213),
        "text_secondary": RGBColor(253, 186, 116),
        "surface": RGBColor(50, 25, 30),
        "chart_bg": "#1e0f14",
        "chart_text": "#ffedd5",
        "chart_colors": ['#fb923c', '#f97316', '#ea580c', '#f59e0b', '#ef4444', '#dc2626'],
        "chart_grid": '#44202a',
    },
}

# ==========================================
# Custom Theme Generator
# ==========================================
CUSTOM_THEME_PRESETS = {
    "blue": {"accent": (59, 130, 246), "bg": (15, 23, 42)},
    "gold": {"accent": (234, 179, 8), "bg": (25, 20, 5)},
    "green": {"accent": (34, 197, 94), "bg": (5, 25, 15)},
    "red": {"accent": (239, 68, 68), "bg": (30, 10, 10)},
    "purple": {"accent": (168, 85, 247), "bg": (20, 5, 35)},
    "pink": {"accent": (236, 72, 153), "bg": (30, 10, 25)},
    "orange": {"accent": (249, 115, 22), "bg": (30, 15, 5)},
    "teal": {"accent": (20, 184, 166), "bg": (5, 25, 25)},
    "cyan": {"accent": (6, 182, 212), "bg": (5, 20, 30)},
    "pastel": {"accent": (186, 186, 255), "bg": (245, 240, 250)},
    "white": {"accent": (30, 64, 175), "bg": (255, 255, 255)},
    "light": {"accent": (30, 64, 175), "bg": (250, 250, 250)},
}

def generate_custom_theme(description: str) -> dict:
    """Generates a theme from a user's text description by matching color keywords."""
    desc_lower = description.lower()
    
    # Try to detect primary and secondary colors
    accent_rgb = (56, 189, 248)  # default accent (sky blue)
    bg_rgb = (15, 23, 42)  # default dark bg
    is_light = False
    
    for keyword, colors in CUSTOM_THEME_PRESETS.items():
        if keyword in desc_lower:
            accent_rgb = colors["accent"]
            bg_rgb = colors["bg"]
            break
    
    # Check for light/dark preference
    if any(w in desc_lower for w in ["light", "white", "bright", "pastel"]):
        is_light = True
        if bg_rgb[0] < 100:  # If bg is still dark, lighten it
            bg_rgb = (245, 245, 250)
    
    if is_light:
        text_rgb = (30, 41, 59)
        text_sec_rgb = (100, 116, 139)
        surface_rgb = (255, 255, 255)
        chart_bg = f"#{bg_rgb[0]:02}{bg_rgb[1]:02}{bg_rgb[2]:02}"
        chart_text = "#1e293b"
        chart_grid = "#e2e8f0"
    else:
        text_rgb = (226, 232, 240)
        text_sec_rgb = (148, 163, 184)
        surface_rgb = (min(bg_rgb[0]+15, 255), min(bg_rgb[1]+15, 255), min(bg_rgb[2]+15, 255))
        chart_bg = f"#{bg_rgb[0]:02}{bg_rgb[1]:02}{bg_rgb[2]:02}"
        chart_text = "#e2e8f0"
        chart_grid = "#334155"
    
    accent_hex = f"#{accent_rgb[0]:02}{accent_rgb[1]:02}{accent_rgb[2]:02}"
    
    return {
        "name": f"Custom: {description[:30]}",
        "bg": RGBColor(*bg_rgb),
        "accent": RGBColor(*accent_rgb),
        "text": RGBColor(*text_rgb),
        "text_secondary": RGBColor(*text_sec_rgb),
        "surface": RGBColor(*surface_rgb),
        "chart_bg": chart_bg,
        "chart_text": chart_text,
        "chart_colors": [accent_hex, '#a855f7', '#ec4899', '#10b981', '#f59e0b', '#ef4444'],
        "chart_grid": chart_grid,
    }


def get_available_themes() -> list:
    """Returns list of available theme metadata for the frontend."""
    return [{"id": tid, "name": t["name"]} for tid, t in PPT_THEMES.items()]

def render_chart_to_image(chart_config: dict, theme: dict = None) -> io.BytesIO:
    """
    Converts Chart.js JSON config to a matplotlib rendered PNG image (BytesIO).
    Supports: bar, line, pie, doughnut, radar
    """
    if theme is None:
        theme = PPT_THEMES["dark"]

    try:
        chart_type = chart_config.get('type', 'bar').lower()
        labels = chart_config.get('labels', [])
        data = chart_config.get('data', [])
        title = chart_config.get('label', 'Chart')
        
        colors = theme['chart_colors']
        
        fig = Figure(figsize=(6, 4), dpi=150)
        fig.patch.set_facecolor(theme['chart_bg'])
        ax = fig.subplots()
        ax.set_facecolor(theme['chart_bg'])
        
        if chart_type == 'bar':
            bars = ax.bar(labels, data, color=colors[:len(labels)], edgecolor='none', width=0.6)
            ax.set_ylabel('Value', color=theme['chart_text'], fontsize=10)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['left'].set_color(theme['chart_grid'])
            ax.spines['bottom'].set_color(theme['chart_grid'])
        elif chart_type == 'line':
            ax.plot(labels, data, marker='o', linewidth=2.5, markersize=8, color=colors[0],
                    markerfacecolor=colors[0], markeredgecolor='white', markeredgewidth=1.5)
            ax.fill_between(range(len(labels)), data, alpha=0.2, color=colors[0])
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['left'].set_color(theme['chart_grid'])
            ax.spines['bottom'].set_color(theme['chart_grid'])
        elif chart_type in ['pie', 'doughnut']:
            wedgeprops = {'width': 0.5, 'edgecolor': theme['chart_bg']} if chart_type == 'doughnut' else {'edgecolor': theme['chart_bg']}
            wedges, texts, autotexts = ax.pie(
                data, labels=labels, autopct='%1.1f%%',
                colors=colors[:len(labels)], wedgeprops=wedgeprops,
                textprops={'color': theme['chart_text']}
            )
            for autotext in autotexts:
                autotext.set_color('white')
                autotext.set_fontweight('bold')
                autotext.set_fontsize(9)
        elif chart_type == 'radar':
            import numpy as np
            angles = [n / len(labels) * 2 * np.pi for n in range(len(labels))]
            data_plot = data + [data[0]]
            angles_plot = angles + [angles[0]]
            ax.remove()
            ax = fig.add_subplot(111, projection='polar')
            ax.set_facecolor(theme['chart_bg'])
            ax.plot(angles_plot, data_plot, 'o-', linewidth=2, color=colors[0])
            ax.fill(angles_plot, data_plot, alpha=0.2, color=colors[0])
            ax.set_xticks(angles)
            ax.set_xticklabels(labels, color=theme['chart_text'], fontsize=9)
            ax.tick_params(axis='y', colors=theme['chart_text'])
            ax.spines['polar'].set_color(theme['chart_grid'])
            ax.grid(color=theme['chart_grid'], alpha=0.5)
        
        # Common styling for non-pie charts
        if chart_type not in ['pie', 'doughnut']:
            ax.tick_params(axis='x', colors=theme['chart_text'], labelsize=9)
            ax.tick_params(axis='y', colors=theme['chart_text'], labelsize=9)
            ax.grid(axis='y', color=theme['chart_grid'], alpha=0.3, linestyle='--')
        
        ax.set_title(title, fontsize=14, fontweight='bold', pad=15, color=theme['chart_text'])
        fig.tight_layout()
        
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                    facecolor=fig.get_facecolor(), edgecolor='none')
        buf.seek(0)
        plt.close(fig)
        
        return buf
    except Exception as e:
        print(f"Chart rendering error: {e}")
        return None


def _add_slide_background(slide, theme):
    """Helper to apply theme background color to a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme['bg']


def _shift_color(rgb: RGBColor, amount: int) -> RGBColor:
    """Lighten (amount>0) or darken (amount<0) an RGBColor, clamped 0-255."""
    r = max(0, min(255, rgb[0] + amount))
    g = max(0, min(255, rgb[1] + amount))
    b = max(0, min(255, rgb[2] + amount))
    return RGBColor(r, g, b)


def _add_gradient_background(slide, theme):
    """Multi-color gradient backdrop (bg → accent-tinted) behind slide content.
    python-pptx's slide background API can't do gradients, so we drop a
    full-bleed rectangle, give it a 2-stop diagonal gradient, and send it to the
    back. Falls back to the solid theme background on any failure."""
    try:
        rect = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        rect.line.fill.background()
        rect.shadow.inherit = False
        rect.fill.gradient()
        stops = rect.fill.gradient_stops
        # Blend bg toward the accent so each deck reads as multi-color, not flat.
        is_light = (theme['bg'][0] + theme['bg'][1] + theme['bg'][2]) > 384
        stops[0].position = 0.0
        stops[0].color.rgb = theme['bg']
        stops[1].position = 1.0
        stops[1].color.rgb = _shift_color(
            RGBColor((theme['bg'][0] + theme['accent'][0]) // 2,
                     (theme['bg'][1] + theme['accent'][1]) // 2,
                     (theme['bg'][2] + theme['accent'][2]) // 2),
            -30 if is_light else 18)
        try:
            rect.fill.gradient_angle = 60.0
        except Exception:
            pass
        # Push the gradient rectangle to the very back so text sits above it.
        spTree = slide.shapes._spTree
        spTree.remove(rect._element)
        spTree.insert(2, rect._element)
        return True
    except Exception as e:
        print(f"PPT gradient fallback (solid bg): {e}")
        _add_slide_background(slide, theme)
        return False


# ---------------------------------------------------------------------------
# Keyless topic-image fetch (DuckDuckGo). Best-effort + cached + bounded; any
# failure (offline, no result, bad bytes) returns None and the slide stays
# image-less. No API key required.
# ---------------------------------------------------------------------------
_IMG_CACHE = {}


def fetch_topic_image(query: str, timeout: int = 6):
    """Return a BytesIO of a topic-relevant JPG/PNG via DuckDuckGo image search,
    or None. Results are cached per-query for the life of the process."""
    if not query or not query.strip():
        return None
    key = query.strip().lower()
    if key in _IMG_CACHE:
        buf = _IMG_CACHE[key]
        if buf is not None:
            buf.seek(0)
            return io.BytesIO(buf.getvalue())
        return None
    try:
        import requests
        headers = {"User-Agent": "Mozilla/5.0 (Cortex PPT)"}
        # 1) get the vqd token DuckDuckGo requires for the image endpoint
        tok = requests.get("https://duckduckgo.com/", params={"q": query},
                           headers=headers, timeout=timeout)
        m = re.search(r'vqd=["\']?([\d-]+)["\']?', tok.text) or \
            re.search(r'vqd=([\d-]+)&', tok.text)
        if not m:
            _IMG_CACHE[key] = None
            return None
        vqd = m.group(1)
        # 2) query the image endpoint
        r = requests.get("https://duckduckgo.com/i.js", headers={**headers,
                         "Referer": "https://duckduckgo.com/"},
                         params={"l": "us-en", "o": "json", "q": query,
                                 "vqd": vqd, "f": ",,,", "p": "1"}, timeout=timeout)
        results = r.json().get("results", [])
        for item in results[:6]:
            url = item.get("image") or ""
            if not url.lower().split("?")[0].endswith((".jpg", ".jpeg", ".png")):
                continue
            try:
                img = requests.get(url, headers=headers, timeout=timeout)
                ctype = img.headers.get("Content-Type", "")
                if img.status_code != 200 or not ctype.startswith("image/"):
                    continue
                if len(img.content) < 1500 or len(img.content) > 5_000_000:
                    continue
                # Magic-byte sanity check (JPEG / PNG)
                if not (img.content[:3] == b"\xff\xd8\xff" or img.content[:8] ==
                        b"\x89PNG\r\n\x1a\n"):
                    continue
                buf = io.BytesIO(img.content)
                _IMG_CACHE[key] = io.BytesIO(img.content)
                buf.seek(0)
                return buf
            except Exception:
                continue
        _IMG_CACHE[key] = None
        return None
    except Exception as e:
        print(f"🖼  [PPT] topic image fetch failed for '{query}': {e}")
        _IMG_CACHE[key] = None
        return None


def _add_decorative_accent_bar(slide, theme):
    """Adds a thin accent-colored bar at the top of content slides."""
    from pptx.util import Inches, Emu
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Emu(54000)  # ~0.06 inches tall
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme['accent']
    shape.line.fill.background()


def _add_slide_number(slide, slide_num, total, theme):
    """Adds a slide number footer."""
    footer_box = slide.shapes.add_textbox(
        Inches(8.5), Inches(7.1), Inches(1.2), Inches(0.3)
    )
    tf = footer_box.text_frame
    tf.text = f"{slide_num} / {total}"
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = theme['text_secondary']
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


# ==========================================
# Intelligent Content Splitter
# ==========================================
def split_narrative_into_sections(narrative: str, target_slides: int) -> list:
    """
    Splits a narrative into sections for slides. Uses ## headings as primary
    delimiters, then intelligently distributes content to hit the target count.
    
    Returns: list of dicts with 'title' and 'bullets' keys
    """
    if not narrative or not narrative.strip():
        return []
    
    # Extract any <CHART> blocks first
    chart_pattern = r'<CHART>(.*?)</CHART>'
    charts_raw = re.findall(chart_pattern, narrative, re.DOTALL)
    # Remove chart blocks from narrative for text processing
    clean_narrative = re.sub(chart_pattern, '', narrative, flags=re.DOTALL).strip()
    # Belt-and-suspenders: if the model's output was truncated mid-chart, a
    # dangling "<CHART>{..." with no closing tag would otherwise dump raw JSON
    # onto a slide. Drop any dangling opener (to end-of-string) and stray tags.
    clean_narrative = re.sub(r'<CHART>.*$', '', clean_narrative, flags=re.DOTALL)
    clean_narrative = clean_narrative.replace('</CHART>', '').replace('<CHART>', '').strip()
    
    # Try splitting by ## headings
    heading_pattern = r'(?:^|\n)\s*#{1,3}\s+(.+)'
    parts = re.split(heading_pattern, clean_narrative)
    
    sections = []
    
    if len(parts) > 1:
        # We have headings — pair each heading with its content
        # parts[0] is text before first heading (intro), parts[1] is first heading, parts[2] is first content, etc.
        if parts[0].strip():
            # There's intro text before first heading
            sections.append({
                'title': 'Overview',
                'bullets': _extract_bullets(parts[0])
            })
        
        for i in range(1, len(parts), 2):
            heading = parts[i].strip().strip('*').strip()
            content = parts[i+1].strip() if i+1 < len(parts) else ''
            bullets = _extract_bullets(content)
            if heading or bullets:
                sections.append({
                    'title': heading,
                    'bullets': bullets
                })
    else:
        # No headings — split by double newlines or bullet points
        lines = [l.strip() for l in clean_narrative.split('\n') if l.strip()]
        
        # Group lines into sections of ~5 bullets each
        chunk_size = max(3, len(lines) // max(target_slides - 1, 1))
        for i in range(0, len(lines), chunk_size):
            chunk = lines[i:i+chunk_size]
            if chunk:
                # Use first line as title if it looks like one
                first = chunk[0]
                if len(first) < 80 and not first.startswith(('•', '-', '*', '→')):
                    sections.append({
                        'title': first.strip('*').strip(':').strip(),
                        'bullets': [_clean_bullet(l) for l in chunk[1:] if l.strip()]
                    })
                else:
                    sections.append({
                        'title': f'Section {len(sections) + 1}',
                        'bullets': [_clean_bullet(l) for l in chunk if l.strip()]
                    })
    
    # If we have fewer sections than target, try to split larger sections
    while len(sections) < target_slides - 1 and any(len(s['bullets']) > 4 for s in sections):
        new_sections = []
        for s in sections:
            if len(s['bullets']) > 4 and len(new_sections) + (len(sections) - len(new_sections)) < target_slides - 1:
                mid = len(s['bullets']) // 2
                new_sections.append({
                    'title': s['title'],
                    'bullets': s['bullets'][:mid]
                })
                new_sections.append({
                    'title': f"{s['title']} (cont.)",
                    'bullets': s['bullets'][mid:]
                })
            else:
                new_sections.append(s)
        if len(new_sections) == len(sections):
            break  # No more splits possible
        sections = new_sections
    
    # If we have more sections than target, merge small ones
    while len(sections) > target_slides - 1 and len(sections) > 1:
        # Find smallest adjacent pair to merge
        min_size = float('inf')
        merge_idx = 0
        for i in range(len(sections) - 1):
            combined = len(sections[i]['bullets']) + len(sections[i+1]['bullets'])
            if combined < min_size:
                min_size = combined
                merge_idx = i
        
        sections[merge_idx]['bullets'].extend(sections[merge_idx + 1]['bullets'])
        sections.pop(merge_idx + 1)
    
    return sections


def _extract_bullets(text: str) -> list:
    """Extract bullet points from a block of text."""
    lines = text.strip().split('\n')
    bullets = []
    for line in lines:
        cleaned = line.strip()
        if cleaned:
            bullets.append(_clean_bullet(cleaned))
    return bullets


def _clean_bullet(text: str) -> str:
    """Clean up a single bullet point text."""
    # Remove markdown bold markers
    text = text.strip()
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    # Remove leading bullet characters
    text = re.sub(r'^[\•\-\*\→]\s*', '', text)
    # Remove leading numbers like "1." or "1)"
    text = re.sub(r'^\d+[\.\)]\s*', '', text)
    return text.strip()


def create_ppt_from_session(session_data: dict) -> io.BytesIO:
    """
    Main PPT generator with theme support, multi-slide splitting, 
    and custom theme generation.
    """
    theme_id = session_data.get('theme', 'dark')
    custom_theme_desc = session_data.get('customTheme', '')
    target_slide_count = session_data.get('slideCount', 10)
    
    # Resolve theme
    if theme_id == 'custom' and custom_theme_desc:
        theme = generate_custom_theme(custom_theme_desc)
    else:
        theme = PPT_THEMES.get(theme_id, PPT_THEMES['dark'])
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    messages = session_data.get('messages', [])
    
    # Combine all narrative content from messages
    full_narrative = ""
    all_charts = []
    
    for msg in messages:
        narrative = msg.get('narrative', '')
        charts = msg.get('charts', [])
        if narrative:
            full_narrative += narrative + "\n\n"
        all_charts.extend(charts)
    
    # Extract charts from the narrative text as well
    chart_pattern = r'<CHART>(.*?)</CHART>'
    chart_matches = re.findall(chart_pattern, full_narrative, re.DOTALL)
    for chart_json in chart_matches:
        try:
            chart_data = json.loads(chart_json.strip())
            all_charts.append(chart_data)
        except:
            pass
    
    # Split narrative into sections for slides
    sections = split_narrative_into_sections(full_narrative, target_slide_count)
    
    # Calculate total slides: 1 title + content sections + chart slides
    chart_slide_count = (len(all_charts) + 1) // 2 if all_charts else 0  # 2 charts per slide
    total_slides = 1 + len(sections) + chart_slide_count
    
    # ==================== TITLE SLIDE ====================
    title_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(title_slide_layout)
    _add_gradient_background(slide, theme)

    # Decorative accent line
    accent_line = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1), Inches(3.4),
        Inches(8), Emu(36000)  # Thin horizontal line
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = theme['accent']
    accent_line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.8))
    title_frame = title_box.text_frame
    main_title = sections[0]['title'] if sections and sections[0].get('title') else session_data.get('sessionName', 'Analysis Report')
    title_frame.text = main_title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = theme['accent']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = session_data.get('projectName', 'Cortex')
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = theme['text']
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Metadata footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(1))
    footer_frame = footer_box.text_frame
    timestamp = datetime.now().strftime("%B %d, %Y at %H:%M")
    model = session_data.get('modelUsed', 'Unknown Model')
    footer_frame.text = f"Generated on {timestamp}"
    footer_frame.paragraphs[0].font.size = Pt(11)
    footer_frame.paragraphs[0].font.color.rgb = theme['text_secondary']
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Model info line
    p2 = footer_frame.add_paragraph()
    p2.text = f"Model: {model}  •  Theme: {theme['name']}"
    p2.font.size = Pt(10)
    p2.font.color.rgb = theme['text_secondary']
    p2.alignment = PP_ALIGN.CENTER
    
    # ==================== CONTENT SLIDES ====================
    slide_num = 1
    
    # Bound how many slides trigger a network image fetch so generation never
    # hangs; deeper slides stay text-only.
    _img_budget = 4

    for section in sections:
        slide_num += 1

        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        _add_gradient_background(slide, theme)
        _add_decorative_accent_bar(slide, theme)
        _add_slide_number(slide, slide_num, total_slides, theme)

        # Topic-relevant photo (right column) for the first few content slides.
        slide_image = None
        if _img_budget > 0 and section.get('title'):
            slide_image = fetch_topic_image(section['title'])
            _img_budget -= 1
        if slide_image is not None:
            try:
                slide.shapes.add_picture(slide_image, Inches(6.5), Inches(1.7),
                                         width=Inches(3.0), height=Inches(2.25))
            except Exception as _img_err:
                print(f"PPT add_picture skipped: {_img_err}")
                slide_image = None

        # Bullets get a narrower column when an image is present so they don't
        # collide with it; otherwise use the full width.
        _bullet_w = Inches(5.6) if slide_image is not None else Inches(8.8)

        # Section title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.text = section['title']
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = theme['accent']
        title_frame.paragraphs[0].space_after = Pt(8)
        
        # Thin divider line under title
        divider = slide.shapes.add_shape(
            1, Inches(0.6), Inches(1.3),
            Inches(3), Emu(18000)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = theme['accent']
        divider.line.fill.background()
        
        # Bullet points
        if section['bullets']:
            text_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), _bullet_w, Inches(5.2))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            
            for idx, bullet in enumerate(section['bullets'][:10]):  # Ma 10 bullets per slide
                if idx == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"• {bullet}" if not bullet.startswith(('•', '-', '→')) else bullet
                p.font.size = Pt(14)
                p.font.color.rgb = theme['text']
                p.space_before = Pt(6)
                p.space_after = Pt(6)
                p.level = 0
    
    # ==================== CHART SLIDES ====================
    if all_charts:
        # Place up to 2 charts per slide
        for i in range(0, len(all_charts), 2):
            slide_num += 1
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            _add_gradient_background(slide, theme)
            _add_decorative_accent_bar(slide, theme)
            _add_slide_number(slide, slide_num, total_slides, theme)

            charts_on_slide = all_charts[i:i+2]
            
            if len(charts_on_slide) == 1:
                # Single centered chart
                img_buf = render_chart_to_image(charts_on_slide[0], theme)
                if img_buf:
                    slide.shapes.add_picture(img_buf, Inches(1.5), Inches(1.5),
                                            width=Inches(7), height=Inches(5))
            else:
                # Two charts side by side
                for j, chart in enumerate(charts_on_slide):
                    img_buf = render_chart_to_image(chart, theme)
                    if img_buf:
                        chart_left = Inches(0.3 + j * 5)
                        slide.shapes.add_picture(img_buf, chart_left, Inches(1.5),
                                                width=Inches(4.5), height=Inches(3.5))
    
    # ==================== THANK YOU SLIDE ====================
    slide_num += 1
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _add_gradient_background(slide, theme)

    # Thank you text
    ty_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    ty_frame = ty_box.text_frame
    ty_frame.text = "Thank You"
    ty_frame.paragraphs[0].font.size = Pt(48)
    ty_frame.paragraphs[0].font.bold = True
    ty_frame.paragraphs[0].font.color.rgb = theme['accent']
    ty_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle (removed "Generated by Cortex")
    
    # Accent bar at bottom
    shape = slide.shapes.add_shape(
        1, Inches(3), Inches(4.8),
        Inches(4), Emu(36000)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme['accent']
    shape.line.fill.background()
    
    # Save to BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output