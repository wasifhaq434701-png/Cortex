import os
import re
import csv
import sqlite3
import datetime
import pandas as pd

# --- Graceful Optional Imports ---
# These are only needed for specific file types. If not installed, 
# the relevant extractor will return empty string instead of crashing.
try:
    from pypdf import PdfReader
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False
    print("⚠ pypdf not installed — PDF parsing disabled. Run: pip install pypdf")

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("⚠ python-docx not installed — DOCX parsing disabled. Run: pip install python-docx")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("⚠ python-pptx not installed — PPTX parsing disabled. Run: pip install python-pptx")

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False
    print("⚠ openpyxl not installed — XLSX parsing disabled. Run: pip install openpyxl")

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    import pytesseract
    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False
    print("⚠ pytesseract not installed — Image OCR disabled. Run: pip install pytesseract")

def generate_tabular_profile(file_path: str) -> str:
    """Dynamically reads CSV, Excel, DB, or raw SQL scripts and generates a statistical profile."""
    ext = file_path.lower().split('.')[-1]
    summary = []
    
    try:
        if ext in ['csv', 'xlsx']:
            df = pd.read_csv(file_path) if ext == 'csv' else pd.read_excel(file_path)
            summary.append(f"--- DATASET PROFILE: {os.path.basename(file_path)} ---")
            summary.append(f"Total Rows: {len(df)} | Total Columns: {len(df.columns)}")
            
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                summary.append("\nNumeric Column Statistics (Mean, Min, Max):")
                summary.append(df[numeric_cols].agg(['mean', 'min', 'max']).round(2).to_string())
                
            cat_cols = df.select_dtypes(include=['object', 'category']).columns
            if len(cat_cols) > 0:
                summary.append("\nCategorical Top Values:")
                for col in cat_cols:
                    summary.append(f"- {col}: {df[col].value_counts().head(5).to_dict()}")

        elif ext in ['db', 'sqlite', 'sqlite3', 'sql']:
            # THE GHOST DATABASE: If it's a raw .sql file, run it in RAM first!
            if ext == 'sql':
                conn = sqlite3.connect(':memory:')
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    conn.executescript(f.read())
            else:
                conn = sqlite3.connect(file_path)
                
            tables = pd.read_sql_query("SELECT name FROM sqlite_master WHERE type='table';", conn)
            summary.append(f"--- DATABASE PROFILE: {os.path.basename(file_path)} ---")
            
            for table in tables['name']:
                df = pd.read_sql_query(f"SELECT * FROM {table}", conn)
                summary.append(f"\nTable '{table}' (Rows: {len(df)}):")
                summary.append(f"Columns: {list(df.columns)}")
                
                num_cols = df.select_dtypes(include=['number']).columns
                if len(num_cols) > 0:
                    summary.append("\nNumeric Averages:")
                    summary.append(df[num_cols].mean().round(2).to_string())
            conn.close()

    except Exception as e:
        return f"Data profiling failed for {os.path.basename(file_path)}: {str(e)}"
        
    return "\n".join(summary)


def clean_text(text: str) -> str:
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def chunk_text(text: str, max_words: int = 400) -> list[str]:
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = []
    current_word_count = 0

    for sentence in sentences:
        words = sentence.split()
        word_count = len(words)
        if current_word_count + word_count > max_words:
            if current_chunk: chunks.append(" ".join(current_chunk))
            current_chunk = [sentence]
            current_word_count = word_count
        else:
            current_chunk.append(sentence)
            current_word_count += word_count
    if current_chunk: chunks.append(" ".join(current_chunk))
    return chunks

# --- Etended File Type Extractors ---
def extract_text_from_pdf(file_path: str) -> str:
    if not HAS_PYPDF:
        print(f"Skipping PDF {file_path} — pypdf not installed.")
        return ""
    text_content = []
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            ext = page.extract_text()
            if ext: text_content.append(ext)
    except Exception as e: print(f"Error parsing PDF {file_path}: {e}")
    return " ".join(text_content)

def extract_text_from_docx(file_path: str) -> str:
    if not HAS_DOCX:
        print(f"Skipping DOCX {file_path} — python-docx not installed.")
        return ""
    try:
        docx = docx.Document(file_path)
        return " ".join([para.text for para in docx.paragraphs if para.text.strip()])
    except Exception: return ""

def extract_text_from_pptx(file_path: str) -> str:
    if not HAS_PPTX:
        print(f"Skipping PPTX {file_path} — python-pptx not installed.")
        return ""
    text_content = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip(): text_content.append(shape.text)
    except Exception: return ""
    return " ".join(text_content)

def extract_text_from_xlsx(file_path: str) -> str:
    text_content = []
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) for cell in row if cell is not None]
                if row_text: text_content.append(" ".join(row_text))
    except Exception: return ""
    return "\n".join(text_content)

def extract_text_from_image(file_path: str) -> str:
    """Extracts text from images using OCR (pytesseract) or returns a descriptive placeholder."""
    if HAS_TESSERACT and HAS_PIL:
        try:
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            if text.strip():
                return f"[IMAGE OCR CONTENT FROM {os.path.basename(file_path)}]\n{text.strip()}"
        except Exception as e:
            print(f"OCR failed for {file_path}: {e}")
    
    # Fallback: Return image metadata even without OCR
    if HAS_PIL:
        try:
            img = Image.open(file_path)
            w, h = img.size
            mode = img.mode
            fmt = img.format or 'Unknown'
            return f"[IMAGE FILE: {os.path.basename(file_path)} | Format: {fmt} | Size: {w}{h} | Mode: {mode}]"
        except Exception:
            pass
    
    return f"[IMAGE FILE: {os.path.basename(file_path)} — OCR not available]"

def extract_text_from_csv(file_path: str, delimiter=',') -> str:
    text_content = []
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            for row in csv.reader(f, delimiter=delimiter): text_content.append(" ".join(row))
    except Exception: return ""
    return "\n".join(text_content)

def extract_metadata_from_db(file_path: str) -> str:
    """Connects to SQLite database files and extracts schemas, tables, and data context."""
    db_summary = []
    try:
        conn = sqlite3.connect(file_path)
        cursor = conn.cursor()
        
        # Fetch all master tables
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = cursor.fetchall()
        
        for table_tuple in tables:
            table_name = table_tuple[0]
            db_summary.append(f"--- Table: {table_name} ---")
            
            # Fetch table schema details
            cursor.execute(f"PRAGMA table_info({table_name});")
            columns = cursor.fetchall()
            col_details = [f"{col[1]} ({col[2]})" for col in columns]
            db_summary.append(f"Columns: {', '.join(col_details)}")
            
            # Fetch a sample snippet of rows safely for context awareness
            cursor.execute(f"SELECT * FROM {table_name} LIMIT 3;")
            sample_rows = cursor.fetchall()
            if sample_rows:
                db_summary.append("Sample Records Context:")
                for row in sample_rows:
                    db_summary.append(str(row))
        conn.close()
    except Exception as e:
        return f"Database introspection error on {os.path.basename(file_path)}: {str(e)}"
    return "\n".join(db_summary)

def get_all_files(path: str) -> list[str]:
    """Returns a flat list of all absolute file paths within a directory (or a single file)."""
    if not os.path.exists(path):
        return []
    
    IGNORE_DIRS = {
        '.git', 'venv', '.venv', 'node_modules', '__pycache__', '.idea', 'build', 'dist',
        '.dart_tool', '.flutter-plugins', '.flutter-plugins-dependencies',
        '.symlinks', 'ephemeral', '.pub-cache', '.pub',
        'Pods', 'DerivedData', '.cworkspace', 'cuserdata',
        '.gradle', '.c',
        'egg-info', '.eggs', '.to', '.mypy_cache', '.pytest_cache',
        '.cache', '.npm', '.yarn', 'coverage', '.next', '.nut',
        'vendor', 'bower_components', '.svn', '.hg',
    }
    
    all_files = []
    if os.path.isfile(path):
        all_files.append(os.path.abspath(path))
    else:
        for root, dirs, files in os.walk(path):
            dirs[:] = [d for d in dirs if d not in IGNORE_DIRS 
                       and not d.endswith(('.codeproj', '.cworkspace'))]
            for file in files:
                all_files.append(os.path.abspath(os.path.join(root, file)))
    
    return all_files

def process_directory(path: str) -> list[dict]:
    """Legacy wrapper for process_files that performs a full scan without incremental caching."""
    all_files = get_all_files(path)
    print(f"📁 Found {len(all_files)} candidate files for full scan")
    return process_files(all_files)

def process_files(file_paths: list[str]) -> list[dict]:
    """Processes a specific list of files, bypassing full directory scans for incremental RAG."""
    processed_data = []
    
    CODE_EXTS = (
        '.py', '.c', '.cpp', '.h', '.hpp', '.java', '.js', '.ts', '.js', '.ts', 
        '.go', '.rs', '.rb', '.php', '.sh', '.sql', '.html', '.css', '.json', 
        '.yaml', '.yml', '.ml', '.md', '.tt', '.ini', '.cfg', '.conf', '.dart'
    )
    
    for file_path in file_paths:
        if not os.path.exists(file_path):
            continue
            
        file = os.path.basename(file_path)
        raw_content = ""
        
        # Determine processing path by extension mapping
        if file.lower().endswith(CODE_EXTS):
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f: 
                    raw_content = f.read()
            except Exception: pass
        elif file.lower().endswith('.pdf'): raw_content = extract_text_from_pdf(file_path)
        elif file.lower().endswith('.docx'): raw_content = extract_text_from_docx(file_path)
        elif file.lower().endswith('.pptx'): raw_content = extract_text_from_pptx(file_path)
        elif file.lower().endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp')):
            raw_content = extract_text_from_image(file_path)
        elif file.lower().endswith(('.csv', '.xlsx', '.db', '.sqlite', '.sqlite3', '.sql')): 
                
                # Get the raw text (for searchability)
                if file.lower().endswith(('.csv', '.xlsx')):
                    raw_content = extract_text_from_xlsx(file_path) if file.endswith('.xlsx') else extract_text_from_csv(file_path)
                else:
                    # For .db or .sql, just grab the raw schema or code
                    try:
                        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                            raw_content = f.read()[:5000] # Cap raw text to save AI memory
                    except Exception: 
                        raw_content = "Binary database file."
                
                # Prepend the Mathematical Profile!
                profile = generate_tabular_profile(file_path)
                raw_content = f"{profile}\n\n{raw_content}"

        if raw_content.strip():
            mtime = os.path.getmtime(file_path)
            date_str = datetime.datetime.fromtimestamp(mtime).strftime('%B %d, %Y')
            
            cleaned = clean_text(raw_content)
            chunks = chunk_text(cleaned)
            
            import hashlib
            path_hash = hashlib.md5(file_path.encode()).hexdigest()[:8]
            
            for index, chunk in enumerate(chunks):
                agent_chunk = f"[SYSTEM LOG - File: {file} | Modified: {date_str}]\n{chunk}"
                processed_data.append({
                    "id": f"{file}_{path_hash}_{index}",
                    "source_file": file_path,
                    "content": agent_chunk
                })
                    
    return processed_data